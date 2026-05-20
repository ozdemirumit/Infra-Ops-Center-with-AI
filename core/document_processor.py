"""
Document Processing Module.
Reads PDF, DOCX, TXT, MD, PPTX files and splits them into chunks.
Prepares text for the RAG pipeline.
"""

import os
from pathlib import Path
from logging_config.logger import get_logger

logger = get_logger("document_processor")

# ─── Supported File Types ───
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".pptx", ".html"}


def extract_text(file_path: str, max_pages: int = None, max_chars: int = None, progress_cb=None) -> str:
    """
    Extracts text content based on file extension.

    Args:
        file_path: Path to the file
        max_pages: For PDFs/PPTX, max pages/slides to read
        max_chars: Stop early once extracted text exceeds this size
        progress_cb: Optional callable(current, total) for progress UI

    Returns:
        Plain text extracted from the file
    """
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(file_path, max_pages=max_pages, progress_cb=progress_cb, max_chars=max_chars)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext == ".html":
        return _extract_html_file(file_path)
    elif ext in (".txt", ".md"):
        return _extract_text_file(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list[dict]:
    """
    Splits text into chunks of a specified size.

    Args:
        text: Text to split
        chunk_size: Maximum character count per chunk
        overlap: Overlap amount between chunks

    Returns:
        [{"index": 0, "text": "..."}, ...]
    """
    if not text or not text.strip():
        return []

    chunks = []
    start = 0
    index = 0

    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]

        # Split at word boundary (last space)
        if end < len(text):
            last_space = chunk.rfind(" ")
            if last_space > chunk_size // 2:
                end = start + last_space
                chunk = text[start:end]

        chunk = chunk.strip()
        if chunk:
            chunks.append({"index": index, "text": chunk})
            index += 1

        start = end - overlap if end < len(text) else len(text)

    return chunks


def process_document(file_path: str, chunk_size: int = 1000, overlap: int = 200) -> list[dict]:
    """
    Read document + chunk it.

    Returns:
        [{"index": 0, "text": "chunk text"}, ...]
    """
    logger.info(f"Processing document: {file_path}")
    text = extract_text(file_path)

    if not text.strip():
        logger.warning(f"Document is empty or unreadable: {file_path}")
        return []

    chunks = chunk_text(text, chunk_size, overlap)
    logger.info(f"Document chunked: {len(chunks)} chunks created ({Path(file_path).name})")
    return chunks


def extract_from_url(url: str, save_dir: str) -> tuple[str, str]:
    """
    Extracts text from a web URL and saves it as HTML.

    Args:
        url: Web page URL
        save_dir: Directory path for saving

    Returns:
        (saved_file_path, extracted_text)
    """
    import requests
    from bs4 import BeautifulSoup
    from urllib.parse import urlparse
    import hashlib

    logger.info(f"Fetching URL content: {url}")

    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
    }
    response = requests.get(url, headers=headers, timeout=30)
    response.raise_for_status()

    soup = BeautifulSoup(response.text, "html.parser")

    # Remove unnecessary elements
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    # Page title
    title = soup.title.string.strip() if soup.title and soup.title.string else ""

    # Extract text
    text = soup.get_text(separator="\n", strip=True)

    # Clean empty lines
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    clean_text = "\n".join(lines)

    if title:
        clean_text = f"# {title}\n\n{clean_text}"

    # Generate filename
    parsed = urlparse(url)
    domain = parsed.netloc.replace("www.", "")
    path_part = parsed.path.strip("/").replace("/", "_")[:50]
    url_hash = hashlib.md5(url.encode()).hexdigest()[:8]
    safe_name = f"{domain}_{path_part}_{url_hash}.html" if path_part else f"{domain}_{url_hash}.html"
    safe_name = "".join(c if c.isalnum() or c in "._-" else "_" for c in safe_name)

    # Save HTML
    save_path = Path(save_dir) / safe_name
    save_path.write_text(response.text, encoding="utf-8")

    logger.info(f"URL saved: {save_path} ({len(clean_text)} characters)")
    return str(save_path), clean_text


# ─── File Readers ───

def _extract_pdf(file_path: str, max_pages: int = None, progress_cb=None, max_chars: int = None) -> str:
    """
    Extracts text from a PDF file.

    Args:
        file_path: Path to the PDF
        max_pages: Optional maximum number of pages to process
        progress_cb: Optional callable(current, total) called per page
        max_chars: Stop extraction once total text exceeds this size
    """
    from PyPDF2 import PdfReader

    reader = PdfReader(file_path)
    total_pages = len(reader.pages)
    limit = min(total_pages, max_pages) if max_pages else total_pages

    pages = []
    total_chars = 0
    for i in range(limit):
        try:
            text = reader.pages[i].extract_text()
        except Exception:
            text = ""
        if text:
            pages.append(f"[Page {i + 1}]\n{text}")
            total_chars += len(text)

        if progress_cb:
            try:
                progress_cb(i + 1, limit)
            except Exception:
                pass

        # Early exit when reaching size cap (no point extracting more — AI truncates anyway)
        if max_chars and total_chars >= max_chars:
            pages.append(f"\n[Truncated at page {i + 1} / {total_pages} — reached size limit]")
            break

    return "\n\n".join(pages)


def _extract_docx(file_path: str) -> str:
    """Extracts text from a DOCX file."""
    from docx import Document

    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    return "\n\n".join(paragraphs)


def _extract_pptx(file_path: str) -> str:
    """Extracts text from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        if texts:
            slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_text_file(file_path: str) -> str:
    """Extracts text from TXT/MD files (auto-detects encoding)."""
    import chardet

    with open(file_path, "rb") as f:
        raw = f.read()

    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8") or "utf-8"

    return raw.decode(encoding, errors="replace")


def _extract_html_file(file_path: str) -> str:
    """Extracts text from an HTML file."""
    from bs4 import BeautifulSoup

    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        html = f.read()

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    title = soup.title.string.strip() if soup.title and soup.title.string else ""
    text = soup.get_text(separator="\n", strip=True)
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    clean_text = "\n".join(lines)

    if title:
        clean_text = f"# {title}\n\n{clean_text}"

    return clean_text
